import io
import os
from dotenv import load_dotenv
from werkzeug.datastructures import FileStorage

# Load environment variables and set up debug functionality
load_dotenv()
DEBUG = os.getenv("DEBUG", "False").lower() in ["true", "1", "t", "yes", "y"]

def debug_log(message):
    if DEBUG:
        print(f"DEBUG [PowerPoint Processor]: {message}")

def process_powerpoint(file_obj):
    """
    Process a PowerPoint presentation (.pptx) and extract its text content.
    
    Args:
        file_obj: File object from request.files
        
    Returns:
        str: Extracted text content in markdown format
    """
    debug_log(f"Processing PowerPoint file: {getattr(file_obj, 'filename', 'Unnamed Presentation')}")
    
    try:
        # Import pptx package here to avoid dependency issues if not installed
        try:
            import pptx
        except ImportError:
            debug_log("python-pptx package not installed. Installing...")
            import subprocess
            subprocess.check_call(["pip", "install", "python-pptx"])
            import pptx
            debug_log("python-pptx package installed successfully")
        
        # Get the file content
        if isinstance(file_obj, FileStorage):
            pptx_bytes = file_obj.read()
            file_obj.close()
            pptx_stream = io.BytesIO(pptx_bytes)
        else:
            # If it's already a bytes object or BytesIO
            pptx_stream = file_obj if isinstance(file_obj, io.BytesIO) else io.BytesIO(file_obj)
        
        # Open the PowerPoint presentation
        presentation = pptx.Presentation(pptx_stream)
        debug_log(f"PowerPoint presentation opened successfully with {len(presentation.slides)} slides")
        
        # Initialize presentation content
        content = f"# Presentation: {getattr(file_obj, 'filename', 'Unnamed Presentation')}\n\n"
        
        # Extract text from each slide
        for i, slide in enumerate(presentation.slides):
            slide_num = i + 1
            content += f"## Slide {slide_num}\n\n"
            
            # Get slide title
            if slide.shapes.title and slide.shapes.title.text:
                content += f"### {slide.shapes.title.text.strip()}\n\n"
            
            # Extract text from all shapes in the slide
            text_parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Skip title text as it's already included
                    if shape == slide.shapes.title:
                        continue
                    text_parts.append(shape.text.strip())
            
            # Add bullet points for text elements
            if text_parts:
                for text in text_parts:
                    # Check if it looks like a bullet point already
                    lines = text.split('\n')
                    for line in lines:
                        line = line.strip()
                        if not line:
                            continue
                        # Add bullet point if it doesn't start with one
                        if not line.startswith('•') and not line.startswith('-'):
                            content += f"- {line}\n"
                        else:
                            content += f"{line}\n"
                content += "\n"
            
            # Extract any tables in the slide
            tables = [shape for shape in slide.shapes if shape.has_table]
            for j, table_shape in enumerate(tables):
                table = table_shape.table
                content += f"\n**Table {j+1} on Slide {slide_num}**\n\n"
                
                # Create markdown table
                rows = []
                for row_idx in range(len(table.rows)):
                    row_cells = []
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        cell_text = cell.text.strip() if cell.text else ""
                        row_cells.append(cell_text)
                    rows.append(row_cells)
                
                if rows:
                    # Create header row
                    header_row = rows[0]
                    content += "| " + " | ".join(header_row) + " |\n"
                    content += "| " + " | ".join(["---" for _ in header_row]) + " |\n"
                    
                    # Add data rows
                    for row in rows[1:]:
                        # Pad shorter rows to match header length
                        row = row + [""] * (len(header_row) - len(row))
                        content += "| " + " | ".join(row) + " |\n"
                    
                    content += "\n"
                
            content += "\n"
        
        debug_log(f"PowerPoint processing complete: {len(content)} characters extracted")
        return content.strip()
    except Exception as e:
        debug_log(f"ERROR processing PowerPoint: {str(e)}")
        return f"Error processing PowerPoint presentation: {str(e)}"
